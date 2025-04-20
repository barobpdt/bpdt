import json
import os
from datetime import datetime
from pathlib import Path
from typing import Union, List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import pandas as pd
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side

def add_image_to_slide(slide, image_path: Union[str, Path], left: float = 1, top: float = 1, 
                      width: float = 4, height: float = 3) -> bool:
    """
    슬라이드에 이미지를 추가하는 함수
    
    Args:
        slide: PowerPoint 슬라이드 객체
        image_path (Union[str, Path]): 이미지 파일 경로
        left (float): 이미지의 왼쪽 위치 (인치)
        top (float): 이미지의 상단 위치 (인치)
        width (float): 이미지의 너비 (인치)
        height (float): 이미지의 높이 (인치)
        
    Returns:
        bool: 이미지 추가 성공 여부
    """
    try:
        # 이미지 경로를 Path 객체로 변환
        img_path = Path(image_path) if isinstance(image_path, str) else image_path
        
        # 이미지 파일이 존재하는지 확인
        if not img_path.exists():
            print(f"이미지 파일이 존재하지 않습니다: {img_path}")
            return False
        
        # 이미지 추가
        slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), 
                               width=Inches(width), height=Inches(height))
        return True
    except Exception as e:
        print(f"이미지 추가 중 오류 발생: {str(e)}")
        return False

def json_to_pptx(json_file: Union[str, Path], output_file: Optional[Union[str, Path]] = None, 
                template_file: Optional[Union[str, Path]] = None) -> tuple[bool, str]:
    """
    JSON 파일의 배열 정보를 PowerPoint 슬라이드로 변환하는 함수
    
    Args:
        json_file (Union[str, Path]): JSON 파일 경로
        output_file (Optional[Union[str, Path]]): 출력 PowerPoint 파일 경로 (기본값: JSON 파일과 같은 이름, 확장자만 .pptx)
        template_file (Optional[Union[str, Path]]): 템플릿 PowerPoint 파일 경로 (기본값: None, 빈 프레젠테이션 사용)
        
    Returns:
        tuple[bool, str]: (성공 여부, 메시지)
    """
    try:
        # JSON 파일 경로를 Path 객체로 변환
        json_path = Path(json_file) if isinstance(json_file, str) else json_file
        
        # JSON 파일이 존재하는지 확인
        if not json_path.exists():
            return False, f"JSON 파일이 존재하지 않습니다: {json_path}"
        
        # JSON 파일 읽기
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # 데이터가 리스트인지 확인
        if not isinstance(data, list):
            return False, f"JSON 데이터는 배열 형식이어야 합니다: {json_path}"
        
        # 출력 파일 경로 설정
        if output_file is None:
            output_path = json_path.with_suffix('.pptx')
        else:
            output_path = Path(output_file) if isinstance(output_file, str) else output_file
        
        # 템플릿 파일이 지정된 경우 해당 파일을 사용, 아니면 빈 프레젠테이션 생성
        if template_file and Path(template_file).exists():
            prs = Presentation(template_file)
        else:
            prs = Presentation()
        
        # 각 JSON 항목에 대해 슬라이드 생성
        for i, item in enumerate(data):
            # 슬라이드 추가 (제목 및 내용 레이아웃)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # 제목 설정
            title = slide.shapes.title
            if title:
                title.text = f"항목 {i+1}"
            
            # 내용 설정
            content = slide.placeholders[1]
            if content:
                # JSON 항목을 문자열로 변환
                if isinstance(item, dict):
                    # 딕셔너리인 경우 키-값 쌍으로 표시
                    text_content = ""
                    image_path = None
                    
                    for key, value in item.items():
                        # 이미지 경로 필드 확인 (image, img, image_path, img_path 등)
                        if key.lower() in ['image', 'img', 'image_path', 'img_path', 'photo', 'picture'] and isinstance(value, str):
                            image_path = value
                        else:
                            text_content += f"{key}: {value}\n"
                    
                    content.text = text_content
                    
                    # 이미지가 있는 경우 추가
                    if image_path:
                        # 이미지 추가 (텍스트 아래에 배치)
                        add_image_to_slide(slide, image_path, left=1, top=2.5, width=4, height=3)
                else:
                    # 그 외의 경우 문자열로 변환
                    content.text = str(item)
        
        # 프레젠테이션 저장
        prs.save(output_path)
        
        return True, f"PowerPoint 파일이 성공적으로 생성되었습니다: {output_path}"
    
    except Exception as e:
        return False, f"PowerPoint 생성 중 오류 발생: {str(e)}"

def create_table_slide(prs: Presentation, data: List[List[Any]], title: str = "테이블") -> None:
    """
    테이블이 있는 슬라이드를 생성하는 함수
    
    Args:
        prs (Presentation): PowerPoint 프레젠테이션 객체
        data (List[List[Any]]): 테이블 데이터 (2차원 배열)
        title (str): 슬라이드 제목 (기본값: "테이블")
    """
    # 테이블 슬라이드 추가
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # 제목 설정
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title
    
    # 테이블 생성
    rows, cols = len(data), len(data[0]) if data else 0
    left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(0.5 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 테이블 스타일 설정
    table.style = 'Medium Grid 2 Accent 1'
    
    # 데이터 채우기
    for i, row in enumerate(data):
        for j, cell_value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_value)
            
            # 셀 텍스트 정렬
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            
            # 헤더 행 스타일 설정
            if i == 0:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(12)

def json_to_table_pptx(json_file: Union[str, Path], output_file: Optional[Union[str, Path]] = None, 
                      template_file: Optional[Union[str, Path]] = None) -> tuple[bool, str]:
    """
    JSON 파일의 배열 정보를 테이블이 있는 PowerPoint 슬라이드로 변환하는 함수
    
    Args:
        json_file (Union[str, Path]): JSON 파일 경로
        output_file (Optional[Union[str, Path]]): 출력 PowerPoint 파일 경로 (기본값: JSON 파일과 같은 이름, 확장자만 .pptx)
        template_file (Optional[Union[str, Path]]): 템플릿 PowerPoint 파일 경로 (기본값: None, 빈 프레젠테이션 사용)
        
    Returns:
        tuple[bool, str]: (성공 여부, 메시지)
    """
    try:
        # JSON 파일 경로를 Path 객체로 변환
        json_path = Path(json_file) if isinstance(json_file, str) else json_file
        
        # JSON 파일이 존재하는지 확인
        if not json_path.exists():
            return False, f"JSON 파일이 존재하지 않습니다: {json_path}"
        
        # JSON 파일 읽기
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # 데이터가 리스트인지 확인
        if not isinstance(data, list):
            return False, f"JSON 데이터는 배열 형식이어야 합니다: {json_path}"
        
        # 출력 파일 경로 설정
        if output_file is None:
            output_path = json_path.with_suffix('.pptx')
        else:
            output_path = Path(output_file) if isinstance(output_file, str) else output_file
        
        # 템플릿 파일이 지정된 경우 해당 파일을 사용, 아니면 빈 프레젠테이션 생성
        if template_file and Path(template_file).exists():
            prs = Presentation(template_file)
        else:
            prs = Presentation()
        
        # 제목 슬라이드 추가
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        if title:
            title.text = f"{json_path.stem} 데이터"
        if subtitle:
            subtitle.text = f"생성일: {datetime.now().strftime('%Y-%m-%d')}"
        
        # 데이터가 딕셔너리 리스트인 경우 테이블로 변환
        if data and isinstance(data[0], dict):
            # 헤더 행 생성
            headers = list(data[0].keys())
            table_data = [headers]
            
            # 데이터 행 생성
            for item in data:
                row = [item.get(header, "") for header in headers]
                table_data.append(row)
            
            # 테이블 슬라이드 생성
            create_table_slide(prs, table_data, f"{json_path.stem} 테이블")
            
            # 이미지가 있는 경우 각 항목에 대한 이미지 슬라이드 생성
            for i, item in enumerate(data):
                # 이미지 경로 필드 확인
                image_path = None
                for key, value in item.items():
                    if key.lower() in ['image', 'img', 'image_path', 'img_path', 'photo', 'picture'] and isinstance(value, str):
                        image_path = value
                        break
                
                if image_path:
                    # 이미지 슬라이드 생성
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    
                    # 제목 설정
                    title = slide.shapes.title
                    if title:
                        title.text = f"{item.get('name', f'항목 {i+1}')} 이미지"
                    
                    # 이미지 추가
                    add_image_to_slide(slide, image_path, left=1, top=1.5, width=8, height=5)
        else:
            # 일반 배열인 경우 각 항목을 별도 슬라이드로 생성
            for i, item in enumerate(data):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                title = slide.shapes.title
                if title:
                    title.text = f"항목 {i+1}"
                
                content = slide.placeholders[1]
                if content:
                    content.text = str(item)
        
        # 프레젠테이션 저장
        prs.save(output_path)
        
        return True, f"PowerPoint 파일이 성공적으로 생성되었습니다: {output_path}"
    
    except Exception as e:
        return False, f"PowerPoint 생성 중 오류 발생: {str(e)}"

def json_to_image_pptx(json_file: Union[str, Path], output_file: Optional[Union[str, Path]] = None, 
                      template_file: Optional[Union[str, Path]] = None) -> tuple[bool, str]:
    """
    JSON 파일의 배열 정보를 이미지가 포함된 PowerPoint 슬라이드로 변환하는 함수
    
    Args:
        json_file (Union[str, Path]): JSON 파일 경로
        output_file (Optional[Union[str, Path]]): 출력 PowerPoint 파일 경로 (기본값: JSON 파일과 같은 이름, 확장자만 .pptx)
        template_file (Optional[Union[str, Path]]): 템플릿 PowerPoint 파일 경로 (기본값: None, 빈 프레젠테이션 사용)
        
    Returns:
        tuple[bool, str]: (성공 여부, 메시지)
    """
    try:
        # JSON 파일 경로를 Path 객체로 변환
        json_path = Path(json_file) if isinstance(json_file, str) else json_file
        
        # JSON 파일이 존재하는지 확인
        if not json_path.exists():
            return False, f"JSON 파일이 존재하지 않습니다: {json_path}"
        
        # JSON 파일 읽기
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # 데이터가 리스트인지 확인
        if not isinstance(data, list):
            return False, f"JSON 데이터는 배열 형식이어야 합니다: {json_path}"
        
        # 출력 파일 경로 설정
        if output_file is None:
            output_path = json_path.with_suffix('.pptx')
        else:
            output_path = Path(output_file) if isinstance(output_file, str) else output_file
        
        # 템플릿 파일이 지정된 경우 해당 파일을 사용, 아니면 빈 프레젠테이션 생성
        if template_file and Path(template_file).exists():
            prs = Presentation(template_file)
        else:
            prs = Presentation()
        
        # 제목 슬라이드 추가
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        if title:
            title.text = f"{json_path.stem} 이미지 갤러리"
        if subtitle:
            subtitle.text = f"생성일: {datetime.now().strftime('%Y-%m-%d')}"
        
        # 각 JSON 항목에 대해 이미지 슬라이드 생성
        for i, item in enumerate(data):
            # 이미지 경로 필드 확인
            image_path = None
            title_text = f"항목 {i+1}"
            
            if isinstance(item, dict):
                for key, value in item.items():
                    if key.lower() in ['image', 'img', 'image_path', 'img_path', 'photo', 'picture'] and isinstance(value, str):
                        image_path = value
                    elif key.lower() in ['name', 'title', 'caption'] and isinstance(value, str):
                        title_text = value
            
            # 이미지가 있는 경우에만 슬라이드 생성
            if image_path:
                # 이미지 슬라이드 생성
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                
                # 제목 설정
                title = slide.shapes.title
                if title:
                    title.text = title_text
                
                # 이미지 추가
                add_image_to_slide(slide, image_path, left=1, top=1.5, width=8, height=5)
                
                # 설명 추가 (있는 경우)
                if isinstance(item, dict) and 'description' in item:
                    left, top, width, height = Inches(1), Inches(6.5), Inches(8), Inches(1)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = item['description']
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 프레젠테이션 저장
        prs.save(output_path)
        
        return True, f"PowerPoint 파일이 성공적으로 생성되었습니다: {output_path}"
    
    except Exception as e:
        return False, f"PowerPoint 생성 중 오류 발생: {str(e)}"

def extract_tables_from_pptx(pptx_file: Union[str, Path], output_file: Optional[Union[str, Path]] = None) -> tuple[bool, str]:
    """
    PowerPoint 파일에서 테이블을 추출하여 Excel 파일로 저장하는 함수
    
    Args:
        pptx_file (Union[str, Path]): PowerPoint 파일 경로
        output_file (Optional[Union[str, Path]]): 출력 Excel 파일 경로 (기본값: PowerPoint 파일과 같은 이름, 확장자만 .xlsx)
        
    Returns:
        tuple[bool, str]: (성공 여부, 메시지)
    """
    try:
        # PowerPoint 파일 경로를 Path 객체로 변환
        pptx_path = Path(pptx_file) if isinstance(pptx_file, str) else pptx_file
        
        # PowerPoint 파일이 존재하는지 확인
        if not pptx_path.exists():
            return False, f"PowerPoint 파일이 존재하지 않습니다: {pptx_path}"
        
        # 출력 파일 경로 설정
        if output_file is None:
            output_path = pptx_path.with_suffix('.xlsx')
        else:
            output_path = Path(output_file) if isinstance(output_file, str) else output_file
        
        # PowerPoint 파일 열기
        prs = Presentation(pptx_path)
        
        # 테이블 데이터를 저장할 리스트
        tables_data = []
        
        # 각 슬라이드에서 테이블 추출
        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if shape.has_table:
                    table = shape.table
                    
                    # 테이블 데이터 추출
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data.append(row_data)
                    
                    # 테이블 데이터가 비어있지 않은 경우에만 추가
                    if table_data and any(any(cell for cell in row) for row in table_data):
                        # 슬라이드 제목 추출 (있는 경우)
                        slide_title = "제목 없음"
                        for title_shape in slide.shapes:
                            if title_shape.has_text_frame and title_shape.text.strip():
                                slide_title = title_shape.text.strip()
                                break
                        
                        # 테이블 정보 저장
                        tables_data.append({
                            "slide_index": slide_index + 1,
                            "shape_index": shape_index + 1,
                            "slide_title": slide_title,
                            "data": table_data
                        })
        
        # 테이블이 없는 경우
        if not tables_data:
            return False, f"PowerPoint 파일에서 테이블을 찾을 수 없습니다: {pptx_path}"
        
        # Excel 파일 생성
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            # 각 테이블을 별도의 시트로 저장
            for i, table_info in enumerate(tables_data):
                # DataFrame 생성
                df = pd.DataFrame(table_info["data"])
                
                # 시트 이름 설정 (최대 31자)
                sheet_name = f"테이블_{i+1}_{table_info['slide_title'][:20]}"
                sheet_name = sheet_name.replace('/', '_').replace('\\', '_')
                
                # DataFrame을 Excel 시트로 저장
                df.to_excel(writer, sheet_name=sheet_name, index=False, header=False)
                
                # 시트 스타일 설정
                worksheet = writer.sheets[sheet_name]
                
                # 헤더 스타일 설정
                header_font = Font(bold=True, size=12)
                header_fill = PatternFill(start_color="E0E0E0", end_color="E0E0E0", fill_type="solid")
                border = Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                )
                
                # 첫 번째 행을 헤더로 설정
                for col in range(1, len(table_info["data"][0]) + 1):
                    cell = worksheet.cell(row=1, column=col)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.border = border
                    cell.alignment = Alignment(horizontal='center', vertical='center')
                
                # 모든 셀에 테두리 및 정렬 설정
                for row in range(1, len(table_info["data"]) + 1):
                    for col in range(1, len(table_info["data"][0]) + 1):
                        cell = worksheet.cell(row=row, column=col)
                        cell.border = border
                        cell.alignment = Alignment(horizontal='center', vertical='center')
                
                # 열 너비 자동 조정
                for col in range(1, len(table_info["data"][0]) + 1):
                    max_length = 0
                    column = worksheet[openpyxl.utils.get_column_letter(col)]
                    for cell in column:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(str(cell.value))
                        except:
                            pass
                    adjusted_width = (max_length + 2)
                    worksheet.column_dimensions[openpyxl.utils.get_column_letter(col)].width = adjusted_width
        
        return True, f"Excel 파일이 성공적으로 생성되었습니다: {output_path}"
    
    except Exception as e:
        return False, f"Excel 파일 생성 중 오류 발생: {str(e)}"

def extract_all_tables_from_pptx(pptx_file: Union[str, Path], output_file: Optional[Union[str, Path]] = None) -> tuple[bool, str]:
    """
    PowerPoint 파일에서 모든 테이블을 추출하여 하나의 Excel 파일에 저장하는 함수
    
    Args:
        pptx_file (Union[str, Path]): PowerPoint 파일 경로
        output_file (Optional[Union[str, Path]]): 출력 Excel 파일 경로 (기본값: PowerPoint 파일과 같은 이름, 확장자만 .xlsx)
        
    Returns:
        tuple[bool, str]: (성공 여부, 메시지)
    """
    try:
        # PowerPoint 파일 경로를 Path 객체로 변환
        pptx_path = Path(pptx_file) if isinstance(pptx_file, str) else pptx_file
        
        # PowerPoint 파일이 존재하는지 확인
        if not pptx_path.exists():
            return False, f"PowerPoint 파일이 존재하지 않습니다: {pptx_path}"
        
        # 출력 파일 경로 설정
        if output_file is None:
            output_path = pptx_path.with_suffix('.xlsx')
        else:
            output_path = Path(output_file) if isinstance(output_file, str) else output_file
        
        # PowerPoint 파일 열기
        prs = Presentation(pptx_path)
        
        # 모든 테이블 데이터를 저장할 리스트
        all_tables_data = []
        
        # 각 슬라이드에서 테이블 추출
        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if shape.has_table:
                    table = shape.table
                    
                    # 테이블 데이터 추출
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data.append(row_data)
                    
                    # 테이블 데이터가 비어있지 않은 경우에만 추가
                    if table_data and any(any(cell for cell in row) for row in table_data):
                        # 슬라이드 제목 추출 (있는 경우)
                        slide_title = "제목 없음"
                        for title_shape in slide.shapes:
                            if title_shape.has_text_frame and title_shape.text.strip():
                                slide_title = title_shape.text.strip()
                                break
                        
                        # 테이블 정보 저장
                        all_tables_data.append({
                            "slide_index": slide_index + 1,
                            "shape_index": shape_index + 1,
                            "slide_title": slide_title,
                            "data": table_data
                        })
        
        # 테이블이 없는 경우
        if not all_tables_data:
            return False, f"PowerPoint 파일에서 테이블을 찾을 수 없습니다: {pptx_path}"
        
        # Excel 파일 생성
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            # 모든 테이블을 하나의 시트에 저장
            combined_data = []
            
            for table_info in all_tables_data:
                # 테이블 구분을 위한 빈 행 추가
                if combined_data:
                    combined_data.append([""] * max(len(row) for row in combined_data))
                
                # 슬라이드 정보 추가
                slide_info = f"슬라이드 {table_info['slide_index']}: {table_info['slide_title']}"
                combined_data.append([slide_info] + [""] * (len(table_info["data"][0]) - 1))
                
                # 테이블 데이터 추가
                for row in table_info["data"]:
                    combined_data.append(row)
            
            # DataFrame 생성
            df = pd.DataFrame(combined_data)
            
            # DataFrame을 Excel 시트로 저장
            df.to_excel(writer, sheet_name="모든 테이블", index=False, header=False)
            
            # 시트 스타일 설정
            worksheet = writer.sheets["모든 테이블"]
            
            # 스타일 설정
            header_font = Font(bold=True, size=12)
            header_fill = PatternFill(start_color="E0E0E0", end_color="E0E0E0", fill_type="solid")
            slide_info_font = Font(bold=True, italic=True)
            slide_info_fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # 스타일 적용
            for row_idx, row in enumerate(combined_data, 1):
                for col_idx, cell_value in enumerate(row, 1):
                    cell = worksheet.cell(row=row_idx, column=col_idx)
                    
                    # 셀 값이 있는 경우에만 스타일 적용
                    if cell_value:
                        # 슬라이드 정보 행 스타일
                        if row_idx > 1 and col_idx == 1 and combined_data[row_idx-2][0] == "":
                            cell.font = slide_info_font
                            cell.fill = slide_info_fill
                        # 테이블 헤더 행 스타일 (슬라이드 정보 다음 행)
                        elif row_idx > 1 and col_idx == 1 and combined_data[row_idx-2][0] != "" and combined_data[row_idx-2][0] != "":
                            cell.font = header_font
                            cell.fill = header_fill
                        
                        cell.border = border
                        cell.alignment = Alignment(horizontal='center', vertical='center')
            
            # 열 너비 자동 조정
            for col in range(1, df.shape[1] + 1):
                max_length = 0
                column = worksheet[openpyxl.utils.get_column_letter(col)]
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = (max_length + 2)
                worksheet.column_dimensions[openpyxl.utils.get_column_letter(col)].width = adjusted_width
        
        return True, f"Excel 파일이 성공적으로 생성되었습니다: {output_path}"
    
    except Exception as e:
        return False, f"Excel 파일 생성 중 오류 발생: {str(e)}"

# 사용 예시
if __name__ == "__main__":
    # 샘플 JSON 파일 생성
    sample_data = [
        {"id": 1, "name": "홍길동", "age": 30, "city": "서울", "image": "sample_images/person1.jpg"},
        {"id": 2, "name": "김철수", "age": 25, "city": "부산", "image": "sample_images/person2.jpg"},
        {"id": 3, "name": "이영희", "age": 28, "city": "인천", "image": "sample_images/person3.jpg"},
        {"id": 4, "name": "박민수", "age": 35, "city": "대구", "image": "sample_images/person4.jpg"},
        {"id": 5, "name": "정지영", "age": 22, "city": "광주", "image": "sample_images/person5.jpg"}
    ]
    
    # 샘플 이미지 디렉토리 생성
    os.makedirs("sample_images", exist_ok=True)
    
    # 샘플 이미지 생성 (실제 이미지가 없는 경우를 위한 예시)
    for i in range(1, 6):
        img = Image.new('RGB', (300, 300), color=(i*50, i*40, i*30))
        img.save(f"sample_images/person{i}.jpg")
    
    # JSON 파일 저장
    with open("sample.json", "w", encoding="utf-8") as f:
        json.dump(sample_data, f, ensure_ascii=False, indent=2)
    
    # JSON을 PowerPoint로 변환
    success, message = json_to_pptx("sample.json")
    print(message)
    
    # JSON을 테이블이 있는 PowerPoint로 변환
    success, message = json_to_table_pptx("sample.json", "sample_table.pptx")
    print(message)
    
    # JSON을 이미지가 포함된 PowerPoint로 변환
    success, message = json_to_image_pptx("sample.json", "sample_image.pptx")
    print(message)
    
    # PowerPoint에서 테이블 추출하여 Excel로 저장
    success, message = extract_tables_from_pptx("sample_table.pptx", "sample_table.xlsx")
    print(message)
    
    # PowerPoint에서 모든 테이블을 하나의 Excel 파일로 저장
    success, message = extract_all_tables_from_pptx("sample_table.pptx", "sample_all_tables.xlsx")
    print(message) 
import pandas as pd
import requests
import json
import matplotlib.pyplot as plt
import yfinance as yf
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches
import os
import matplotlib
matplotlib.use('Agg')  # GUI 없이 실행하기 위한 설정

def get_kospi_listed_companies():
    """
    주가 등락률이 큰 상위 5개 종목의 월간 데이터를 가져와서 그래프를 그리고 파워포인트로 저장하는 코드를 작성하겠습니다:    
    KRX에서 코스피 상장회사 정보를 가져오는 함수
    """
    # KRX 요청 URL
    url = "http://data.krx.co.kr/comm/bldAttendant/getJsonData.cmd"
    
    # 요청 파라미터 - 최신 KRX API 구조에 맞게 수정
    params = {
        "bld": "dbms/MDC/STAT/standard/MDCSTAT03901",
        "locale": "ko_KR",
        "mktId": "STK",  # 주식시장
        "trdDd": datetime.now().strftime("%Y%m%d"),  # 오늘 날짜
        "money": "1",  # 원화
        "csvxls_isNo": "false",  # JSON 형식으로 요청
    }
    
    # 요청 헤더
    headers = {
        "Content-Type": "application/x-www-form-urlencoded; charset=UTF-8",
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36",
        "Accept": "application/json, text/javascript, */*; q=0.01",
        "Accept-Language": "ko-KR,ko;q=0.9,en-US;q=0.8,en;q=0.7",
        "Origin": "http://data.krx.co.kr",
        "Referer": "http://data.krx.co.kr/contents/MDC/MDCSTAT03901/index.jsp",
    }
    
    try:
        # 데이터 요청
        response = requests.post(url, data=params, headers=headers)
        response.raise_for_status()  # 오류 발생 시 예외 발생
        
        # JSON 데이터 파싱
        data = json.loads(response.text)
        
        # 데이터프레임으로 변환
        if 'OutBlock_1' in data:
            df = pd.DataFrame(data['OutBlock_1'])
            
            # 컬럼명 한글로 변경
            column_mapping = {
                'ISU_ABBRV': '종목명',
                'ISU_SHRT_CD': '종목코드',
                'ISU_NM': '종목전체명',
                'MKT_NM': '시장구분',
                'SEC_NM': '섹터',
                'TDD_CLSPRC': '종가',
                'CMPPREVDD_PRC': '전일대비',
                'FLUC_RT': '등락률',
                'TDD_OPNPRC': '시가',
                'TDD_HGPRC': '고가',
                'TDD_LWPRC': '저가',
                'ACC_TRDVOL': '거래량',
                'ACC_TRDVOL_VAL': '거래대금',
                'MKTCAP': '시가총액',
                'LIST_SHRS': '상장주식수',
            }
            
            # 실제 응답에 있는 컬럼만 매핑
            available_columns = {k: v for k, v in column_mapping.items() if k in df.columns}
            df = df.rename(columns=available_columns)
            
            # 숫자형 데이터 변환
            numeric_columns = ['종가', '전일대비', '등락률', '시가', '고가', '저가', '거래량', '거래대금', '시가총액', '상장주식수']
            for col in numeric_columns:
                if col in df.columns:
                    df[col] = pd.to_numeric(df[col].astype(str).str.replace(',', ''), errors='coerce')
            
            return df
        else:
            print("응답 데이터에 'OutBlock_1' 키가 없습니다.")
            return None
    
    except Exception as e:
        print(f"데이터를 가져오는 중 오류가 발생했습니다: {e}")
        return None

def get_monthly_stock_data(stock_code, months=6):
    """
    주어진 종목 코드에 대한 월간 주가 데이터를 가져오는 함수
    """
    try:
        # 종목 코드에 .KS 추가 (코스피 종목)
        ticker = f"{stock_code}.KS"
        
        # yfinance를 사용하여 주가 데이터 가져오기
        stock = yf.Ticker(ticker)
        
        # 현재 날짜로부터 지정된 개월 수 전까지의 데이터 가져오기
        end_date = datetime.now()
        start_date = end_date - timedelta(days=months*30)
        
        # 월간 데이터 가져오기
        hist = stock.history(start=start_date, end=end_date, interval="1mo")
        
        return hist
    except Exception as e:
        print(f"종목 {stock_code}의 데이터를 가져오는 중 오류가 발생했습니다: {e}")
        return None

def create_stock_trend_graph(stock_data, stock_name, output_path):
    """
    주가 추이 그래프를 생성하고 저장하는 함수
    """
    try:
        plt.figure(figsize=(10, 6))
        plt.plot(stock_data.index, stock_data['Close'], marker='o', linewidth=2)
        plt.title(f"{stock_name} 월간 주가 추이", fontsize=14)
        plt.xlabel("날짜", fontsize=12)
        plt.ylabel("주가 (원)", fontsize=12)
        plt.grid(True, linestyle='--', alpha=0.7)
        plt.xticks(rotation=45)
        plt.tight_layout()
        
        # 그래프 저장
        plt.savefig(output_path, dpi=300, bbox_inches='tight')
        plt.close()
        
        return True
    except Exception as e:
        print(f"그래프 생성 중 오류가 발생했습니다: {e}")
        return False

def create_powerpoint_with_graphs(top_stocks, stock_data_dict, output_file="stock_trend_analysis.pptx"):
    """
    그래프를 포함한 파워포인트 프레젠테이션을 생성하는 함수
    """
    try:
        # 새 프레젠테이션 생성
        prs = Presentation()
        
        # 제목 슬라이드 추가
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "코스피 상위 종목 주가 추이 분석"
        subtitle.text = f"생성일: {datetime.now().strftime('%Y-%m-%d')}"
        
        # 각 종목별 슬라이드 추가
        for i, (_, stock) in enumerate(top_stocks.iterrows()):
            stock_code = stock['종목코드']
            stock_name = stock['종목명']
            
            if stock_code in stock_data_dict and stock_data_dict[stock_code] is not None:
                # 그래프 이미지 경로
                graph_path = f"graph_{stock_code}.png"
                
                # 그래프 생성
                create_stock_trend_graph(stock_data_dict[stock_code], stock_name, graph_path)
                
                # 슬라이드 추가
                slide_layout = prs.slide_layouts[5]  # 제목과 내용이 있는 레이아웃
                slide = prs.slides.add_slide(slide_layout)
                
                # 제목 설정
                title = slide.shapes.title
                title.text = f"{stock_name} ({stock_code}) 주가 추이"
                
                # 그래프 추가
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                height = Inches(5)
                slide.shapes.add_picture(graph_path, left, top, width=width, height=height)
                
                # 임시 그래프 파일 삭제
                if os.path.exists(graph_path):
                    os.remove(graph_path)
        
        # 프레젠테이션 저장
        prs.save(output_file)
        print(f"파워포인트 프레젠테이션이 {output_file}에 저장되었습니다.")
        return True
    except Exception as e:
        print(f"파워포인트 생성 중 오류가 발생했습니다: {e}")
        return False

def main():
    # 코스피 상장회사 정보 가져오기
    print("코스피 상장회사 정보를 가져오는 중...")
    df = get_kospi_listed_companies()
    
    if df is not None:
        # 등락률 기준으로 상위 5개 종목 선택
        top_stocks = df.sort_values(by='등락률', ascending=False).head(5)
        print("\n등락률 상위 5개 종목:")
        print(top_stocks[['종목명', '종목코드', '등락률']])
        
        # 각 종목의 월간 데이터 가져오기
        stock_data_dict = {}
        for _, stock in top_stocks.iterrows():
            stock_code = stock['종목코드']
            stock_name = stock['종목명']
            print(f"\n{stock_name}({stock_code})의 월간 데이터를 가져오는 중...")
            
            stock_data = get_monthly_stock_data(stock_code)
            if stock_data is not None and not stock_data.empty:
                stock_data_dict[stock_code] = stock_data
                print(f"{stock_name} 데이터 가져오기 성공: {len(stock_data)}개 데이터")
            else:
                stock_data_dict[stock_code] = None
                print(f"{stock_name} 데이터 가져오기 실패")
        
        # 파워포인트 프레젠테이션 생성
        if stock_data_dict:
            create_powerpoint_with_graphs(top_stocks, stock_data_dict)
        else:
            print("주가 데이터를 가져오지 못했습니다.")
    else:
        print("코스피 상장회사 정보를 가져오지 못했습니다.")

if __name__ == "__main__":
    main() 
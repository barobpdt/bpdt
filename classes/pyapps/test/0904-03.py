
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width=Inches(16)
prs.slide_height=Inches(9)
slide=prs.slides.add_slide(prs.slide_layouts[6])
shapes=slide.shapes
shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(9))
#Title
# shape.title.text = "Title of the slide"
# Shape position
left = Inches(0.5)
top = Inches(1.5)
width = Inches(2.0)
height = Inches(0.2)

box = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#Fill
fill = box.fill
line = box.line
fill.solid()
fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2
line.color.theme_color = MSO_THEME_COLOR.ACCENT_2
# How can I set font, size of text for the shape ?
# One Object for instead two seperate ones
#box_text.font.bold = True

# Text position
t_left = Inches(0.5)
t_top = Inches(1.4)
t_width = Inches(2.0)
t_height = Inches(0.4)
#Text
txBox = slide.shapes.add_textbox(t_left, t_top, t_width,t_height)
tf = txBox.text_frame.paragraphs[0]
tf.vertical_anchor = MSO_ANCHOR.TOP
tf.word_wrap = False
tf.margin_top = 0
tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
run = tf.add_run()
run.text = "Text on the Shape"
font = run.font
font.name = 'Calibri'
font.size = Pt(18)
font.bold = True
font.italic = None  # cause value to be inherited from theme
font.color.theme_color = MSO_THEME_COLOR.ACCENT_5
prs.save('c:/temp/xx/0904-04.pptx')	

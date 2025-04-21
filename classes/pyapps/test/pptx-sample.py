
from pptx import Presentation, oxml
from pptx.oxml import parse_xml
from pptx.table import _Cell
from typing import cast
from copy import deepcopy

fpOut=open("c:/temp/log.txt", 'a') 
def log (msg):
	fpOut.write(f"##> {msg}\n")
	fpOut.flush()

def append_col(tb): 
    new_col = deepcopy(tb._tbl.tblGrid.gridCol_lst[-1])
    tb._tbl.tblGrid.append(new_col)  # copies last grid element
    for tr in tb._tbl.tr_lst:
        new_tc = deepcopy(tr.tc_lst[-1])
        tr.tc_lst[-1].addnext(new_tc)
        cell = _Cell(new_tc, tr.tc_lst)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = '--'
    tblchildren = tb._tbl.getchildren()
    for child in tblchildren:
        if isinstance(child, oxml.table.CT_TableGrid):
            ws = set()
            for j in child:
                if j.w not in ws:
                    ws.add(j.w)
                else:
                    # print('j:\n', j.xml)
                    for elem in j:
                            j.remove(elem)
    return tb
	
pptx_file = r"c:/temp/test.pptx"
table=None
with open(pptx_file, "rb") as pptx:
	prs = Presentation(pptx)
	for slide in prs.slides:
		for shape in slide.shapes:
			if not shape.has_table:
				continue
			table = shape.table

print ("table==>", table, pptx_file)
 
for i, row in enumerate(table.rows):
	# Iterate over the cells of the row
	for j, cell in enumerate(row.cells):
		# Get the table cell properties (tcPr) from the table cell (tc)
		# tcPr = cell._tc.get_or_add_tcPr() 
		log(cell._tc.xml)
		
append_col(table)
prs.save('c:/temp/test-a.ppt')		

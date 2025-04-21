
 
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE    
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

prs = Presentation()
slide=prs.slides.add_slide(prs.slide_layouts[6])
shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(6))
text_frame = shape.text_frame
text_frame.text = 'Text I want to appear in text-box'
text_frame.auto_size = MSO_AUTO_SIZE.NONE
text_frame.word_wrap = False
for paragraph in text_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.LEFT
prs.save('c:/temp/xx/0904-01.pptx')	

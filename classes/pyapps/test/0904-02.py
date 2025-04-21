
 
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE    
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

prs = Presentation('c:/temp/xx/0904-01.pptx')
slide=prs.slides[0]
for shape in slide.shapes:
    print("=======> ", shape._sp, shape.has_text_frame)
    if not shape.has_text_frame:
        continue
        
    # iterate through paragarphs in shape
    for p in shape.text_frame.paragraphs:
        # store formats and their runs by index (not dict because of duplicate runs)
        formats, newRuns = [], []

        # iterate through runs
        for r in p.runs:
            # get text
            text = r.text
            print("text ========== ", text)

            # replace text
            text = text.replace('s','xyz')

            # store run
            newRuns.append(text)

            # store format
            formats.append({'size':r.font.size,
                            'bold':r.font.bold,
                            'underline':r.font.underline,
                            'italic':r.font.italic})

        # clear paragraph
        p.clear()
        print("new run =>", newRuns)

        # iterate through new runs and formats and write to paragraph
        for i in range(len(newRuns)):
            # add run with text
            run = p.add_run()
            run.text = newRuns[i]

            # format run
            run.font.bold = formats[i]['bold']
            run.font.italic = formats[i]['italic']
            run.font.size = formats[i]['size']
            run.font.underline = formats[i]['underline'] 
 
prs.save('c:/temp/xx/0904-02.pptx')	
